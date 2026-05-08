from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

try:
    from pptx import Presentation
    from pptx.util import Inches
except Exception:  # pragma: no cover
    Presentation = None
    Inches = None


ROOT = Path(__file__).resolve().parents[3]
DECK_DIR = ROOT / "docs" / "product-handoff-deck"
OUT = DECK_DIR / "output"
PREVIEWS = DECK_DIR / "previews"
OUT.mkdir(parents=True, exist_ok=True)
PREVIEWS.mkdir(parents=True, exist_ok=True)

W, H = 1600, 900
BG = "#FFFFFF"
INK = "#17212F"
TEXT = "#2E3B4B"
MUTED = "#667180"
BLUE = "#1D5FA7"
BLUE_DARK = "#174C86"
BLUE_PALE = "#EEF5FC"
GRAY_PALE = "#F5F7FA"
LINE = "#CDD6E0"
ORANGE = "#C57918"
GREEN = "#227C55"
RED = "#B83A3A"

FONT = "C:/Windows/Fonts/msyh.ttc"
FONT_BOLD = "C:/Windows/Fonts/msyhbd.ttc"


def font(size, bold=False):
    return ImageFont.truetype(FONT_BOLD if bold else FONT, size)


def lines_for(draw, text, fnt, max_width):
    lines = []
    for para in text.split("\n"):
        cur = ""
        for ch in para:
            trial = cur + ch
            if draw.textlength(trial, font=fnt) <= max_width or not cur:
                cur = trial
            else:
                lines.append(cur)
                cur = ch
        if cur:
            lines.append(cur)
    return lines


def draw_text(draw, xy, text, size=24, fill=TEXT, bold=False, max_width=None, gap=6):
    x, y = xy
    fnt = font(size, bold)
    if max_width is None:
        draw.text((x, y), text, font=fnt, fill=fill)
        return draw.textbbox((x, y), text, font=fnt)[3]
    for line in lines_for(draw, text, fnt, max_width):
        draw.text((x, y), line, font=fnt, fill=fill)
        y += size + gap
    return y


def rect(draw, box, fill=BG, outline=LINE, width=2):
    draw.rectangle(box, fill=fill, outline=outline, width=width)


def line(draw, xy, fill=LINE, width=2):
    draw.line(xy, fill=fill, width=width)


def arrow(draw, x1, y1, x2, y2, color=BLUE):
    draw.line((x1, y1, x2, y2), fill=color, width=4)
    draw.polygon([(x2, y2), (x2 - 15, y2 - 8), (x2 - 15, y2 + 8)], fill=color)


def canvas():
    img = Image.new("RGB", (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, W, 12), fill=BLUE_DARK)
    return img


def header(draw, title, judgement):
    draw_text(draw, (70, 42), title, 41, INK, True, 1460, 4)
    draw.rectangle((70, 108, 150, 114), fill=BLUE)
    draw_text(draw, (70, 132), f"核心判断：{judgement}", 24, TEXT, True, 1430, 5)
    line(draw, (70, 180, 1530, 180), LINE, 2)


def conclusion(draw, text):
    y = 812
    draw.rectangle((70, y - 8, 78, y + 54), fill=ORANGE)
    draw_text(draw, (96, y), "关键结论：", 25, ORANGE, True)
    draw_text(draw, (228, y), text, 25, TEXT, True, 1260, 5)


def footer(draw, idx, total=6):
    draw_text(draw, (70, 862), f"Bank Permission Material AI Pre-audit | {idx}/{total}", 17, "#8A96A3")


def table(draw, x, y, col_widths, row_heights, rows, header_fill=BLUE_PALE, body_size=20):
    cur_y = y
    for r, row_h in enumerate(row_heights):
        cur_x = x
        for c, col_w in enumerate(col_widths):
            fill = header_fill if r == 0 else (GRAY_PALE if r % 2 == 0 else BG)
            rect(draw, (cur_x, cur_y, cur_x + col_w, cur_y + row_h), fill=fill)
            item = rows[r][c]
            color = BLUE if r == 0 else TEXT
            bold = r == 0 or (c == 0 and r > 0)
            size = 22 if r == 0 else body_size
            draw_text(draw, (cur_x + 16, cur_y + 14), item, size, color, bold, col_w - 32, 5)
            cur_x += col_w
        cur_y += row_h


def pill(draw, x, y, text, color=BLUE, fill=BLUE_PALE):
    w = int(draw.textlength(text, font=font(19, True))) + 32
    rect(draw, (x, y, x + w, y + 38), fill=fill, outline=LINE, width=1)
    draw_text(draw, (x + 16, y + 7), text, 19, color, True)
    return w


def save_slide(img, idx):
    path = PREVIEWS / f"slide_{idx:02d}.png"
    img.save(path)
    return path


def slide1():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(
        d,
        "统一定位：构建面向银行材料报送前的 AI 预审能力",
        "项目服务于支付部门会计作业，不替代审批、不形成控制点，核心是降低对外报送错误和复核成本。",
    )

    roles = [
        ("申请人", "准备申请表、附件、证件等材料\n降低漏填、错填、错对象风险"),
        ("审核人", "快速识别一致、待确认和不一致事项\n把复核精力集中在少量关键问题"),
        ("审批人", "关注对外报送前的材料质量和关键风险\n避免陷入字段级技术细节"),
    ]
    x = 100
    for i, (title, body) in enumerate(roles):
        rect(d, (x, 242, x + 430, 438), fill=BLUE_PALE if i == 1 else BG)
        draw_text(d, (x + 26, 274), title, 31, BLUE, True)
        draw_text(d, (x + 26, 334), body, 23, TEXT, False, 365, 7)
        x += 500

    rows = [
        ["业务目标", "当前含义", "对产品设计的约束"],
        ["对外信息准确", "主体、账号、权限、介质、证件等关键信息不能传错", "默认展示业务语言，技术规则码只给产品/IT排查"],
        ["作业秩序稳定", "辅助申请人与审核人聚焦问题，不替代审批责任", "只提示建议复核，不输出阻断、放行或不可提交"],
        ["证据可回看", "每条结论应能回到来源文档和解析片段", "证据预览是报告页的关键能力"],
    ]
    table(d, 110, 510, [260, 540, 540], [54, 74, 74, 74], rows)
    conclusion(d, "这不是 AI 自动审批系统，而是把正式对外报送前的材料预审做成可读、可复核、可迭代的业务能力。")
    footer(d, 1)
    return img


def slide2():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(
        d,
        "拉通主链路：以电子流为基准完成材料解析、映射与检查",
        "电子流是本次办理的基准事实，申请材料围绕它做解析、标准化、比对、归类和证据回看。",
    )
    steps = [
        ("输入材料", "电子流\n申请表 / 附件\n证件图片"),
        ("文档解析", "Word / DOC / PDF\nOCR / Vision"),
        ("要素抽取", "主体 / 账号 / 权限\n介质 / 平台 / 证件"),
        ("检查比对", "硬规则比对\n语义风险检查"),
        ("业务报告", "绿灯通过\n审核人确认\n发现不一致"),
    ]
    x = 75
    for i, (title, body) in enumerate(steps):
        rect(d, (x, 245, x + 255, 420), fill=BLUE_PALE if i in (2, 3) else BG)
        draw_text(d, (x + 20, 276), f"{i + 1}. {title}", 25, BLUE, True)
        draw_text(d, (x + 20, 336), body, 21, TEXT, False, 205, 5)
        if i < len(steps) - 1:
            arrow(d, x + 265, 332, x + 318, 332)
        x += 305

    rows = [
        ["能力环节", "当前 Demo 已覆盖", "产品化仍需补齐"],
        ["多文档解析与抽取", "支持多份 Word/DOC/PDF 与证件图片，输出统一业务结构", "更多真实模板、版式和脱敏样本的回归验证"],
        ["硬规则 + 语义检查", "确定性字段用规则比对，场景和权限理解交给语义检查", "跨文档关系校验、业务规则口径和置信度管理"],
        ["业务友好报告", "互斥统计、分组明细、证据预览、技术信息折叠", "原文页码、坐标、高亮定位以及角色化视图"],
    ]
    table(d, 130, 510, [300, 560, 460], [54, 78, 78, 78], rows)
    conclusion(d, "当前 Demo 已跑通预审主链路；下一阶段重点是用真实样本把模板、字段映射和规则口径资产化。")
    footer(d, 2)
    return img


def slide3():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(
        d,
        "收敛报告阅读：从复杂审计明细转为三类业务结果",
        "会计和财务用户不需要理解 AI 与规则码，只需要知道哪些一致、哪些要人看、哪些建议修正。",
    )
    stats = [
        ("共检查", "N", "项内容", BLUE),
        ("绿灯通过", "X", "一致 / 通过", GREEN),
        ("审核人确认", "Y", "系统无法直接判断", ORANGE),
        ("发现不一致", "Z", "需补充或修正", RED),
    ]
    x = 90
    for title, num, body, color in stats:
        rect(d, (x, 230, x + 320, 375), fill=BG, outline=color, width=4)
        draw_text(d, (x + 28, 255), title, 23, MUTED, True)
        draw_text(d, (x + 28, 300), num, 42, INK, True)
        draw_text(d, (x + 94, 318), body, 19, TEXT, False, 180)
        x += 370

    rows = [
        ["分类", "业务含义", "默认展示方式", "用户动作"],
        ["绿灯通过", "材料与电子流或检查逻辑一致", "折叠展示，避免干扰审核人", "通常无需优先处理"],
        ["审核人确认", "系统当前无法直接判断，或需要结合业务背景", "按组展示少量代表项，其余收起", "查看依据后人工判断"],
        ["发现不一致", "材料与电子流或规则之间存在明确差异", "优先展示并保留证据入口", "补充、修正或重点复核"],
    ]
    table(d, 90, 455, [250, 450, 430, 330], [54, 78, 78, 78], rows)
    conclusion(d, "统计口径采用互斥分类：共检查 = 绿灯通过 + 审核人确认 + 发现不一致，避免业务用户看到数字对不上。")
    footer(d, 3)
    return img


def slide4():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(
        d,
        "重构证据表达：默认给业务依据，技术细节按需折叠",
        "报告应先回答业务问题：电子流登记了什么、材料识别到什么、依据来自哪份文档，而不是先暴露字段名和规则码。",
    )
    rect(d, (120, 235, 740, 690), fill=BLUE_PALE)
    draw_text(d, (150, 265), "业务默认看到", 30, BLUE, True)
    bullets = [
        "电子流中登记的信息",
        "申请材料中识别的信息",
        "来源文档与涉及内容",
        "材料原文 / 解析片段",
        "建议处理方向：确认、补充或修正",
    ]
    y = 336
    for b in bullets:
        draw_text(d, (165, y), f"• {b}", 24, TEXT, False, 520)
        y += 58

    rect(d, (860, 235, 1480, 690), fill=BG)
    draw_text(d, (890, 265), "默认不打扰业务用户", 30, BLUE, True)
    bullets2 = [
        "不默认展示字段簇、规则码、检查方式等内部字段",
        "不默认展示权限布尔值和模型中间结构",
        "不默认展示内部规则编码和调试信息",
        "技术信息统一收进“给产品/IT查看的技术信息”",
    ]
    y = 336
    for b in bullets2:
        draw_text(d, (905, y), f"• {b}", 23, TEXT, False, 520)
        y += 68

    draw_text(d, (690, 438), "→", 56, ORANGE, True)
    conclusion(d, "页面表达的核心变化是从“模型/规则输出”转为“业务可执行提示”：看得懂、知道看哪里、能回到依据。")
    footer(d, 4)
    return img


def slide5():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(
        d,
        "建立运营闭环：反馈进入受控规则资产，而非直接改系统",
        "AI 可以辅助归因和生成候选，但银行材料场景必须经过确认、版本化和回归验证后再生效。",
    )
    top = [("业务反馈", "认可 / 误报 / 漏报\n证据不清 / 字段错误"), ("AI 辅助归因", "判断问题来源\n形成候选变更"), ("候选变更池", "等待确认\n不直接生效")]
    x = 90
    for i, (title, body) in enumerate(top):
        rect(d, (x, 250, x + 330, 400), fill=BLUE_PALE if i == 1 else BG)
        draw_text(d, (x + 25, 278), title, 27, BLUE, True)
        draw_text(d, (x + 25, 332), body, 20, TEXT, False, 260, 4)
        if i < len(top) - 1:
            arrow(d, x + 340, 325, x + 420, 325)
        x += 430

    rows = [
        ["分流类型", "典型内容", "生效路径"],
        ["业务配置类", "字段别名、模板映射、枚举归一、提取优先级", "业务确认 → 配置库更新 → 小样本回归 → 较快生效"],
        ["规则 / Prompt 类", "风险等级、人工确认口径、语义规则补充", "产品复核 → 版本更新 → 回归验证 → 灰度生效"],
        ["工程版本类", "PDF/Word 定位、复杂跨文档校验、解析器增强", "产品需求池 → 工程开发 → 测试验收 → 版本发布"],
    ]
    table(d, 120, 500, [250, 530, 620], [54, 76, 76, 76], rows)
    conclusion(d, "上线后的安全边界是：反馈可沉淀为候选，确认和回归后才进入配置、规则或工程版本。")
    footer(d, 5)
    return img


def slide6():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(
        d,
        "分阶段推进产品化：从 Demo 跑通走向可运营能力",
        "后续重点不是继续堆更多检查项，而是围绕真实样本建设模板、规则、证据链和反馈运营体系。",
    )
    rows = [
        ["阶段", "核心任务", "主要产出", "验收关注"],
        ["阶段一：样本验证", "准备脱敏真实样本，批量跑数，记录误报、漏报和人工确认项", "字段映射表、问题候选表、模板差异表", "业务是否看得懂，风险是否能聚焦"],
        ["阶段二：资产化", "建设模板库、字段映射库、规则库和 Prompt 版本", "可配置规则、回归样本集、发布流程", "哪些可配置快生效，哪些进工程版本"],
        ["阶段三：证据链", "补齐 Word/PDF 原文定位、页码、坐标和高亮", "证据预览、原文定位、审核复核链路", "结论是否可追溯、可解释、可审计"],
        ["阶段四：系统集成", "接入真实电子流、权限、日志、密级和运营看板", "上线方案、运营指标、反馈闭环", "上线门槛和持续运营机制"],
    ]
    table(d, 60, 220, [230, 470, 420, 330], [58, 92, 92, 92, 92], rows, body_size=19)
    conclusion(d, "建议对产品同事明确：当前 Demo 验证方向和主链路，下一阶段交付重点是样本验证、资产化和证据链建设。")
    footer(d, 6)
    return img


def build_pptx(slides):
    if Presentation is None:
        return None
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]
    for path in slides:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    out = OUT / "contract_handle_product_handoff_style_refreshed.pptx"
    prs.save(out)
    return out


def main():
    slide_funcs = [slide1, slide2, slide3, slide4, slide5, slide6]
    paths = []
    for idx, fn in enumerate(slide_funcs, start=1):
        paths.append(save_slide(fn(), idx))
    pptx = build_pptx(paths)
    print("Generated previews:")
    for path in paths:
        print(path)
    if pptx:
        print(f"Generated PPTX: {pptx}")
    else:
        print("python-pptx is unavailable; skipped PPTX export.")


if __name__ == "__main__":
    main()
